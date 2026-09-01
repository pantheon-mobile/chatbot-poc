import ast
import io
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ingestion_powerpoint import _safe_shape_type, parse_pptx_presentation


def _sample_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "奨学金説明"
    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(1))
    text_box.text_frame.text = "第一種は無利子"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(6), Inches(1.5)).table
    table.cell(0, 0).text = "種類"
    table.cell(0, 1).text = "利子"
    table.cell(1, 0).text = "第二種"
    table.cell(1, 1).text = "有利子"
    second = prs.slides.add_slide(prs.slide_layouts[5])
    second.shapes.title.text = "申込手続"
    second.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1)).text = "学校へ相談"
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _load_app_function(name):
    source = Path(__file__).parents[1].joinpath("app.py").read_text(encoding="utf-8")
    tree = ast.parse(source)
    constants = {"INGESTION_TEST_KB_PREFIXES"}
    nodes = [
        node for node in tree.body
        if isinstance(node, ast.FunctionDef) and node.name == name
        or isinstance(node, ast.Assign) and any(
            isinstance(target, ast.Name) and target.id in constants for target in node.targets
        )
    ]
    def build_ingestion_metadata(source_type, ingestion_format, ui_metadata, source_url,
                                 source_file_name, original_title, datasource_id,
                                 original_source_file_name):
        return {"metadataAttributes": {
            "source_type": source_type, "ingestion_format": ingestion_format,
            "source_file_name": source_file_name, "original_title": original_title,
            "datasource_id": datasource_id,
            "original_source_file_name": original_source_file_name
        }}

    namespace = {
        "json": json, "build_ingestion_metadata": build_ingestion_metadata,
        "remove_empty_metadata_attributes": lambda metadata: metadata,
    }
    exec(compile(ast.Module(body=nodes, type_ignores=[]), "app.py", "exec"), namespace)
    return namespace[name]


def test_parse_pptx_preserves_slide_markers_text_and_table():
    result = parse_pptx_presentation(_sample_pptx(), "sample.pptx")
    assert result["slide_count"] == 2
    assert result["table_count"] == 1
    assert "--- slide 1 ---" in result["txt_text"]
    assert "<!-- slide 2 -->" in result["markdown_text"]
    assert "第一種は無利子" in result["markdown_text"]
    assert "| 種類 | 利子 |" in result["markdown_text"]
    assert "| 第二種 | 有利子 |" in result["markdown_text"]


def test_unrecognized_powerpoint_shape_type_is_ignored():
    class UnsupportedShape:
        @property
        def shape_type(self):
            raise NotImplementedError("Shape instance of unrecognized shape type")

    assert _safe_shape_type(UnsupportedShape()) is None


def test_build_ppt_artifacts_has_three_kb_formats_and_sidecars():
    pptx_bytes = _sample_pptx()
    result = parse_pptx_presentation(pptx_bytes, "sample.pptx")
    build = _load_app_function("build_ppt_ingestion_artifacts")
    artifacts, metadata = build(pptx_bytes, "sample.pptx", "DS_TEST", result, {})
    kb_artifacts = [item for item in artifacts if item["role"] == "KB同期用コピー"]
    assert {item["format"] for item in kb_artifacts} == {
        "PPT_PPTX", "PPT_TXT", "PPT_MARKDOWN"
    }
    assert len(kb_artifacts) == 6
    assert metadata["PPT_MARKDOWN"]["metadataAttributes"]["slide_count"] == 2
    assert any(item["key"].startswith(
        "documents/ingestion-test/kb-source/ppt-markdown/DS_TEST/"
    ) for item in kb_artifacts)
