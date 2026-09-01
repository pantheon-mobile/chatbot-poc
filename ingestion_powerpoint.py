import io
import time

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _escape_markdown(value: str) -> str:
    return str(value or "").replace("|", "\\|").replace("\n", "<br>")


def _table_rows(table) -> list[list[str]]:
    return [[cell.text.strip() for cell in row.cells] for row in table.rows]


def _markdown_table(rows: list[list[str]]) -> str:
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    normalized = [row + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    return "\n".join([
        "| " + " | ".join(_escape_markdown(value) for value in header) + " |",
        "| " + " | ".join("---" for _ in range(width)) + " |",
        *("| " + " | ".join(_escape_markdown(value) for value in row) + " |"
          for row in normalized[1:])
    ])


def _safe_shape_type(shape):
    """Return a shape type without aborting on unsupported PowerPoint shapes."""
    try:
        return shape.shape_type
    except (NotImplementedError, ValueError):
        return None


def parse_pptx_presentation(pptx_bytes: bytes, original_name: str = "") -> dict:
    """PPTXをスライド順の構造付きTXTとMarkdownへ変換する。"""
    started = time.perf_counter()
    presentation = Presentation(io.BytesIO(pptx_bytes))
    load_ms = round((time.perf_counter() - started) * 1000, 1)
    txt_sections, markdown_sections, slide_previews = [], [], []
    table_count = image_count = text_shape_count = notes_count = 0
    parse_started = time.perf_counter()

    for slide_number, slide in enumerate(presentation.slides, start=1):
        title = ""
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip()
        txt_lines = [f"--- slide {slide_number} ---"]
        markdown_lines = [f"<!-- slide {slide_number} -->", f"# {title or f'スライド {slide_number}'}"]
        if title:
            txt_lines.append(title)
        preview = {"slide_number": slide_number, "title": title, "texts": [], "tables": [], "notes": ""}

        for shape in slide.shapes:
            # python-pptx can expose some valid Office shapes but raise
            # NotImplementedError when their shape_type is read.  Such a shape
            # must not make the whole presentation conversion fail.
            if _safe_shape_type(shape) == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
            if getattr(shape, "has_table", False):
                rows = _table_rows(shape.table)
                table_count += 1
                preview["tables"].append(rows)
                txt_lines.append("[表]")
                txt_lines.extend("\t".join(row) for row in rows)
                markdown_lines.extend(["", _markdown_table(rows)])
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text or (title and text == title):
                continue
            text_shape_count += 1
            preview["texts"].append(text)
            txt_lines.append(text)
            for paragraph in shape.text_frame.paragraphs:
                paragraph_text = paragraph.text.strip()
                if not paragraph_text:
                    continue
                indent = "  " * int(getattr(paragraph, "level", 0) or 0)
                markdown_lines.append(f"{indent}- {paragraph_text}")

        try:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            notes = ""
        if notes:
            notes_count += 1
            preview["notes"] = notes
            txt_lines.extend(["[発表者ノート]", notes])
            markdown_lines.extend(["", "## 発表者ノート", notes])

        txt_sections.append("\n".join(txt_lines))
        markdown_sections.append("\n".join(markdown_lines))
        slide_previews.append(preview)

    parse_ms = round((time.perf_counter() - parse_started) * 1000, 1)
    txt_started = time.perf_counter()
    txt_text = "\n\n".join(txt_sections).strip()
    txt_generation_ms = round((time.perf_counter() - txt_started) * 1000, 1)
    markdown_started = time.perf_counter()
    markdown_text = "\n\n".join(markdown_sections).strip()
    markdown_generation_ms = round((time.perf_counter() - markdown_started) * 1000, 1)
    title = presentation.core_properties.title or ""
    if not title and slide_previews:
        title = slide_previews[0]["title"]
    return {
        "title": title or original_name, "slide_count": len(presentation.slides),
        "table_count": table_count, "image_count": image_count,
        "text_shape_count": text_shape_count, "notes_count": notes_count,
        "txt_text": txt_text, "markdown_text": markdown_text,
        "slides": slide_previews, "pptx_load_ms": load_ms,
        "slide_parse_ms": parse_ms, "txt_generation_ms": txt_generation_ms,
        "markdown_generation_ms": markdown_generation_ms,
        "generated_file_count": 3,
        "generated_total_bytes": len(pptx_bytes) + len(txt_text.encode("utf-8"))
        + len(markdown_text.encode("utf-8"))
    }
